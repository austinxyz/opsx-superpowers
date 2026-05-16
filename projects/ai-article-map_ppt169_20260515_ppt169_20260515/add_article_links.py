#!/usr/bin/env python3
"""Post-process PPTX to add article hyperlinks to slides 6-10."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

# Per-slide keyword → URL mapping (ordered, first match wins)
SLIDE_LINKS = {
    6: [
        ("引言 Vibe Coding",           "https://juejin.cn/post/7585397173023113235"),
        ("AI 时代的编程新范式",         "https://juejin.cn/post/7585397173023113235"),
        ("Claude Code 概述",            "https://juejin.cn/post/7585412203978047531"),
        ("全流程实战",                   "https://juejin.cn/post/7585397173023178771"),
        ("软件开发方法论在 AI 时代",    "https://juejin.cn/post/7585412203978063915"),
        ("Claude Code 的适用场景",      "https://juejin.cn/post/7585390679398907910"),
        ("项目开发总结",                 "https://juejin.cn/post/7585390679398924294"),
        ("SDD——OpenSpec",               "https://juejin.cn/post/7613975134063591476"),
        ("Claude Code 实战使用心得",    "https://juejin.cn/post/7585400495683502122"),
    ],
    7: [
        ("十倍效率提升",                 "https://juejin.cn/post/7612144065903427636"),
        ("敏捷二十年",                   "https://juejin.cn/post/7614752576312721460"),
        ("OpenSpec × Superpowers",       "https://juejin.cn/post/7630755707556118562"),
        ("三周之后",                     "https://juejin.cn/post/7639268681235333160"),
        ("构建 AI Agent 系统",           "https://juejin.cn/post/7636978905913409576"),
        ("驯服 AI Agent",                "https://juejin.cn/post/7593215107405692964"),
    ],
    8: [
        ("经验越丰富的工程师",           "https://juejin.cn/post/7639239837708730403"),
        ("反而越容易被淘汰",             "https://juejin.cn/post/7639239837708730403"),
        ("Junior 工程师",                "https://juejin.cn/post/7615801634950283298"),
        ("没有未来了吗",                 "https://juejin.cn/post/7615801634950283298"),
        ("软件开发经理",                 "https://juejin.cn/post/7624102246314426422"),
        ("管理杂务缩减到三分之一",       "https://juejin.cn/post/7624102246314426422"),
        ("运维工程师如何",               "https://juejin.cn/post/7618239540529496118"),
        ("Ops → Platform",               "https://juejin.cn/post/7618239540529496118"),
    ],
    9: [
        ("普通人的杠杆",                 "https://juejin.cn/post/7626336868328718376"),
        ("杠杆在哪里",                   "https://juejin.cn/post/7626336868328718376"),
        ("AI 时代，普通人的",            "https://juejin.cn/post/7626336868328718376"),
        ("善用 AI 的三个心法",           "https://juejin.cn/post/7636587578021560355"),
        ("为什么有人越用越好",           "https://juejin.cn/post/7636587578021560355"),
        ("北美理财太复杂",               "https://juejin.cn/post/7635681418753163304"),
        ("Claude Code 搭了知识库",       "https://juejin.cn/post/7635681418753163304"),
        ("龙虾热",                       "https://juejin.cn/post/7618764794955300927"),
        ("OpenClaw",                     "https://juejin.cn/post/7618764794955300927"),
        ("热的是什么",                   "https://juejin.cn/post/7618764794955300927"),
    ],
    10: [
        ("引言 Vibe Coding",             "https://juejin.cn/post/7585397173023113235"),
        ("Claude Code 概述",             "https://juejin.cn/post/7585412203978047531"),
        ("Claude Code 项目开发全流程实战","https://juejin.cn/post/7585397173023178771"),
        ("软件开发方法论在 AI 时代",     "https://juejin.cn/post/7585412203978063915"),
        ("Claude Code 的适用场景",       "https://juejin.cn/post/7585390679398907910"),
        ("Claude Code 项目开发的总结",   "https://juejin.cn/post/7585390679398924294"),
        ("从 Vibe Coding 到 SDD",        "https://juejin.cn/post/7613975134063591476"),
        ("理财系统完整复盘",             "https://juejin.cn/post/7585400495683502122"),
        ("十倍效率提升",                 "https://juejin.cn/post/7612144065903427636"),
        ("敏捷二十年",                   "https://juejin.cn/post/7614752576312721460"),
        ("OpenSpec × Superpowers",       "https://juejin.cn/post/7630755707556118562"),
        ("三周之后：OpenSpec + Superpowers", "https://juejin.cn/post/7639268681235333160"),
        ("构建 AI Agent 系统",           "https://juejin.cn/post/7636978905913409576"),
        ("驯服 AI Agent",                "https://juejin.cn/post/7593215107405692964"),
        ("经验越丰富的工程师反而越容易被淘汰", "https://juejin.cn/post/7639239837708730403"),
        ("Junior 工程师没有未来了吗",    "https://juejin.cn/post/7615801634950283298"),
        ("软件开发经理：管理杂务",       "https://juejin.cn/post/7624102246314426422"),
        ("运维工程师如何不被淘汰",       "https://juejin.cn/post/7618239540529496118"),
        ("普通人的杠杆在哪里",           "https://juejin.cn/post/7626336868328718376"),
        ("善用 AI 的三个心法",           "https://juejin.cn/post/7636587578021560355"),
        ("北美理财太复杂",               "https://juejin.cn/post/7635681418753163304"),
        ("龙虾热",                       "https://juejin.cn/post/7618764794955300927"),
    ],
}


def find_url(slide_num: int, text: str) -> str | None:
    for keyword, url in SLIDE_LINKS.get(slide_num, []):
        if keyword in text:
            return url
    return None


def add_hyperlink(shape, url: str, slide_part, rId_cache: dict) -> None:
    if url not in rId_cache:
        rId_cache[url] = slide_part.relate_to(
            url,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True,
        )
    rId = rId_cache[url]

    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            r_elem = run._r
            rPr = r_elem.find(qn("a:rPr"))
            if rPr is None:
                rPr = etree.SubElement(r_elem, qn("a:rPr"))
                r_elem.insert(0, rPr)

            for el in rPr.findall(qn("a:hlinkClick")):
                rPr.remove(el)

            hlinkClick = etree.SubElement(rPr, qn("a:hlinkClick"))
            hlinkClick.set(qn("r:id"), rId)

            # Prevent hyperlink underline from appearing
            if "u" not in rPr.attrib:
                rPr.set("u", "none")


def main() -> None:
    if len(sys.argv) < 2:
        print("Usage: python add_article_links.py <pptx_file> [output_file]")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    output_path = (
        Path(sys.argv[2]) if len(sys.argv) >= 3
        else pptx_path.with_name(pptx_path.stem + "_linked.pptx")
    )

    prs = Presentation(str(pptx_path))
    total_linked = 0

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_num not in SLIDE_LINKS:
            continue

        rId_cache: dict[str, str] = {}
        slide_linked = 0

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            url = find_url(slide_num, text)
            if not url:
                continue

            add_hyperlink(shape, url, slide.part, rId_cache)
            slide_linked += 1
            print(f"  Slide {slide_num}: '{text[:55].replace(chr(10), ' ')}' -> linked")

        if slide_linked:
            print(f"  -- {slide_linked} links on slide {slide_num}")
        total_linked += slide_linked

    prs.save(str(output_path))
    print(f"\nDone: {total_linked} hyperlinks added")
    print(f"Output: {output_path}")


if __name__ == "__main__":
    main()
